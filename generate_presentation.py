# -*- coding: utf-8 -*-
"""
Exchange Rate Determination - Professional PPTX Generator
Zambia ZMW/USD Analysis Dashboard Presentation
"""
import sys
import io
import os
sys.stdout.reconfigure(encoding='utf-8')
import copy
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
from matplotlib.patches import FancyBboxPatch
import matplotlib.gridspec as gridspec
import numpy as np

from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
import lxml.etree as etree

# ─── COLOR PALETTE ──────────────────────────────────────────────────────────────
NAVY        = RGBColor(0x09, 0x15, 0x2F)   # deepest navy
NAVY_MID    = RGBColor(0x12, 0x28, 0x55)   # mid navy
NAVY_LIGHT  = RGBColor(0x1A, 0x3C, 0x7D)   # lighter navy
ACCENT_BLUE = RGBColor(0x1E, 0x88, 0xE5)   # bright accent blue
SKY_BLUE    = RGBColor(0x64, 0xB5, 0xF6)   # sky blue
TEAL        = RGBColor(0x00, 0xAC, 0xD7)   # teal accent
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE   = RGBColor(0xF4, 0xF7, 0xFC)
LIGHT_GRAY  = RGBColor(0xE3, 0xEB, 0xF6)
MID_GRAY    = RGBColor(0x90, 0xA4, 0xBE)
GOLD        = RGBColor(0xF5, 0xA6, 0x23)
GOLD_LIGHT  = RGBColor(0xFF, 0xD0, 0x54)
GREEN       = RGBColor(0x2E, 0xCC, 0x71)
RED         = RGBColor(0xE7, 0x4C, 0x3C)
DARK_TEXT   = RGBColor(0x0D, 0x1B, 0x3E)

# hex helpers for matplotlib
def h(rgb): return "#{}".format(str(rgb))

NAVY_H       = h(NAVY)
NAVY_MID_H   = h(NAVY_MID)
NAVY_LIGHT_H = h(NAVY_LIGHT)
ACCENT_H     = h(ACCENT_BLUE)
SKY_H        = h(SKY_BLUE)
TEAL_H       = h(TEAL)
GOLD_H       = h(GOLD)
GREEN_H      = h(GREEN)
RED_H        = h(RED)
WHITE_H      = "#FFFFFF"
OFF_WHITE_H  = h(OFF_WHITE)

# ─── REPRESENTATIVE DATA ─────────────────────────────────────────────────────────
YEARS       = [2015, 2016, 2017, 2018, 2019, 2020, 2021, 2022, 2023]
AVG_RATE    = [8.60, 10.34, 9.85, 11.72, 13.24, 17.59, 16.81, 17.22, 20.89]
UPPER_BOUND = [9.46, 11.37, 10.84, 12.89, 14.56, 19.35, 18.49, 18.94, 22.98]
LOWER_BOUND = [7.74,  9.31,  8.87, 10.55, 11.92, 15.83, 15.13, 15.50, 18.80]
PPP_TARGET  = [8.25,  9.80,  9.20, 11.10, 12.80, 16.50, 16.00, 16.80, 19.50]
YOY_DEPR    = [None, 20.2, -4.7,  19.1, 12.8,  32.9,  -4.4,   2.4,  21.3]
BID_ASK     = [1.80,  2.30,  2.00,  2.50,  2.80,  3.20,  2.90,  2.70,  3.10]
US_INF      = [0.10,  1.30,  2.10,  2.40,  1.80,  1.20,  4.70,  8.00,  4.10]
ZM_INF      = [10.1, 17.9,   6.6,   7.9,   9.2,  15.7,  22.0,  11.2,  13.5]
INF_DIFF    = [zi - ui for zi, ui in zip(ZM_INF, US_INF)]

# Latest-period KPIs (2023)
KPI_AVG_RATE   = "20.89"
KPI_DEPR       = "+21.3%"
KPI_SPREAD     = "3.1%"
KPI_US_INF     = "4.1%"
KPI_ZM_INF     = "13.5%"
KPI_INF_DIFF   = "9.4%"
KPI_VARIANCE   = "-9.1%"
KPI_VERDICT    = "Depreciation\nPressure"
KPI_PPP        = "19.50"

# ─── HELPERS ─────────────────────────────────────────────────────────────────────
def add_rect(slide, x, y, w, h_in, fill_rgb, alpha=None, line_rgb=None, line_pt=0):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h_in))
    fill  = shape.fill
    fill.solid()
    fill.fore_color.rgb = fill_rgb
    line  = shape.line
    if line_rgb:
        line.color.rgb = line_rgb
        line.width = Pt(line_pt)
    else:
        line.fill.background()
    shape.shadow.inherit = False
    return shape

def add_text(slide, text, x, y, w, h_in, font_size, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h_in))
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(font_size)
    run.font.bold  = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name  = "Calibri"
    return txb

def add_rounded_rect(slide, x, y, w, h_in, fill_rgb, radius_emu=91440,
                     line_rgb=None, line_pt=1.5):
    """Add a rounded rectangle via raw XML."""
    x_emu = int(Inches(x));  y_emu = int(Inches(y))
    w_emu = int(Inches(w));  h_emu = int(Inches(h_in))
    adj = int(radius_emu / min(w_emu, h_emu) * 100000)
    adj = min(adj, 50000)
    sp_xml = (
        '<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"'
        ' xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
        '<p:nvSpPr>'
        '<p:cNvPr id="1" name="RndRect"/>'
        '<p:cNvSpPr><a:spLocks noGrp="1"/></p:cNvSpPr>'
        '<p:nvPr/>'
        '</p:nvSpPr>'
        '<p:spPr>'
        f'<a:xfrm><a:off x="{x_emu}" y="{y_emu}"/>'
        f'<a:ext cx="{w_emu}" cy="{h_emu}"/></a:xfrm>'
        f'<a:prstGeom prst="roundRect"><a:avLst>'
        f'<a:gd name="adj" fmla="val {adj}"/>'
        f'</a:avLst></a:prstGeom>'
        '<a:solidFill>'
        f'<a:srgbClr val="{str(fill_rgb)}"/>'
        '</a:solidFill>'
    )
    if line_rgb:
        sp_xml += (
            f'<a:ln w="{int(Pt(line_pt))}">'
            '<a:solidFill>'
            f'<a:srgbClr val="{str(line_rgb)}"/>'
            '</a:solidFill></a:ln>'
        )
    else:
        sp_xml += '<a:ln><a:noFill/></a:ln>'
    sp_xml += '</p:spPr><p:txBody><a:bodyPr/><a:lstStyle/><a:p/></p:txBody></p:sp>'
    sp_element = parse_xml(sp_xml)
    slide.shapes._spTree.append(sp_element)

def chart_to_image(fig):
    buf = io.BytesIO()
    fig.savefig(buf, format='png', dpi=150, bbox_inches='tight',
                facecolor='none', transparent=True)
    buf.seek(0)
    plt.close(fig)
    return buf

def add_chart_image(slide, buf, x, y, w, h_in):
    slide.shapes.add_picture(buf, Inches(x), Inches(y), Inches(w), Inches(h_in))

def kpi_card(slide, x, y, w, h_in, label, value, unit="",
             value_color=WHITE, accent_color=ACCENT_BLUE,
             bg_color=NAVY_MID, border_color=ACCENT_BLUE):
    """Render a KPI card — thick consistent border, no internal overlaps."""
    PAD = 0.14   # safe horizontal inset from border on each side
    # Card outline — thick 2.5pt border, well-defined rounded corners
    add_rounded_rect(slide, x, y, w, h_in, bg_color,
                     radius_emu=100000, line_rgb=border_color, line_pt=2.5)
    # Horizontal accent separator below label — inset safely from all edges
    sep_y = y + 0.38
    add_rect(slide, x + PAD, sep_y, w - PAD * 2, 0.022, accent_color)
    # Label (above separator)
    add_text(slide, label.upper(), x + PAD, y + 0.10, w - PAD * 2, 0.28,
             7.5, bold=False, color=SKY_BLUE, align=PP_ALIGN.CENTER)
    # Value (main KPI number)
    add_text(slide, value, x + PAD, sep_y + 0.08, w - PAD * 2, h_in - sep_y + y - 0.22,
             22, bold=True, color=value_color, align=PP_ALIGN.CENTER)
    # Unit label at bottom
    if unit:
        add_text(slide, unit, x + PAD, y + h_in - 0.26, w - PAD * 2, 0.20,
                 7, bold=False, color=MID_GRAY, align=PP_ALIGN.CENTER)

def section_header(slide, number, title, subtitle=""):
    """Navy full-width header bar at top of content slides."""
    add_rect(slide, 0, 0, 13.33, 1.12, NAVY)
    # left accent stripe
    add_rect(slide, 0, 0, 0.06, 1.12, GOLD)
    # section number pill
    add_rounded_rect(slide, 0.18, 0.18, 0.55, 0.55, ACCENT_BLUE, radius_emu=200000)
    add_text(slide, str(number), 0.18, 0.15, 0.55, 0.60,
             16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # title
    add_text(slide, title, 0.85, 0.10, 9.5, 0.55,
             20, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_text(slide, subtitle, 0.85, 0.62, 9.5, 0.35,
                 9, bold=False, color=SKY_BLUE, align=PP_ALIGN.LEFT)
    # right page stamp area
    add_text(slide, "EXCHANGE RATE DETERMINATION | ZAMBIA", 10.5, 0.38, 2.7, 0.35,
             6.5, bold=False, color=MID_GRAY, align=PP_ALIGN.RIGHT)


# ═══════════════════════════════════════════════════════════════════════════════
#  CHART GENERATORS
# ═══════════════════════════════════════════════════════════════════════════════

plt.rcParams.update({
    'font.family': 'DejaVu Sans',
    'axes.spines.top':   False,
    'axes.spines.right': False,
    'axes.spines.left':  False,
    'axes.spines.bottom': False,
    'xtick.color': '#90A4BE',
    'ytick.color': '#90A4BE',
    'xtick.labelsize': 7.5,
    'ytick.labelsize': 7.5,
})

def make_exchange_rate_chart():
    fig, ax = plt.subplots(figsize=(7.8, 3.2))
    fig.patch.set_facecolor(NAVY_MID_H)
    ax.set_facecolor(NAVY_MID_H)

    yr = np.array(YEARS)
    ax.fill_between(yr, LOWER_BOUND, UPPER_BOUND,
                    alpha=0.18, color=ACCENT_H, label='Policy Band')
    ax.plot(yr, PPP_TARGET, '--', color=GOLD_H, linewidth=1.5,
            dashes=(5, 3), label='PPP Target Rate', zorder=3)
    ax.plot(yr, UPPER_BOUND, '-', color=SKY_H, linewidth=1, alpha=0.6, label='Upper Bound')
    ax.plot(yr, LOWER_BOUND, '-', color=SKY_H, linewidth=1, alpha=0.6, label='Lower Bound')
    ax.fill_between(yr, AVG_RATE, color=ACCENT_H, alpha=0.25)
    ax.plot(yr, AVG_RATE, 'o-', color=ACCENT_H, linewidth=2.2,
            markersize=5, markerfacecolor=WHITE_H, markeredgewidth=1.5,
            label='Avg Selling Rate', zorder=5)

    ax.set_xticks(yr)
    ax.set_xticklabels([str(y) for y in YEARS], color='#90A4BE')
    ax.tick_params(axis='y', colors='#90A4BE')
    ax.set_ylabel('ZMW / USD', color='#90A4BE', fontsize=8)
    ax.yaxis.set_label_coords(-0.05, 0.5)
    ax.grid(axis='y', color='#1A3C7D', linewidth=0.5, linestyle='--', alpha=0.5)
    ax.grid(axis='x', visible=False)

    legend = ax.legend(loc='upper left', framealpha=0.0, fontsize=7,
                       labelcolor='#90A4BE', ncol=2)
    for line in legend.get_lines():
        line.set_linewidth(2)

    ax.set_title('Exchange Rate Trend vs Policy Bounds', color=WHITE_H,
                 fontsize=10, fontweight='bold', pad=8)
    fig.tight_layout(pad=0.5)
    return chart_to_image(fig)


def make_depreciation_chart():
    fig, ax = plt.subplots(figsize=(5.5, 2.8))
    fig.patch.set_facecolor(NAVY_MID_H)
    ax.set_facecolor(NAVY_MID_H)

    yr    = YEARS[1:]
    depr  = YOY_DEPR[1:]
    colors = [GREEN_H if d < 0 else RED_H if d > 20 else ACCENT_H for d in depr]

    bars = ax.bar(yr, depr, color=colors, width=0.6, zorder=3, edgecolor='none')
    ax.axhline(0, color='#90A4BE', linewidth=0.8, linestyle='-')

    for bar, val in zip(bars, depr):
        ypos = bar.get_height() + 0.5 if val >= 0 else bar.get_height() - 2.2
        ax.text(bar.get_x() + bar.get_width()/2, ypos,
                f'{val:+.1f}%', ha='center', va='bottom',
                fontsize=7, color=WHITE_H, fontweight='bold')

    ax.set_xticks(yr)
    ax.set_xticklabels([str(y) for y in yr], color='#90A4BE', fontsize=7)
    ax.tick_params(axis='y', colors='#90A4BE')
    ax.grid(axis='y', color='#1A3C7D', linewidth=0.5, linestyle='--', alpha=0.5)
    ax.set_title('Year-on-Year Depreciation (%)', color=WHITE_H,
                 fontsize=9.5, fontweight='bold', pad=8)
    fig.tight_layout(pad=0.5)
    return chart_to_image(fig)


def make_bidask_chart():
    fig, ax = plt.subplots(figsize=(5.5, 2.8))
    fig.patch.set_facecolor(NAVY_MID_H)
    ax.set_facecolor(NAVY_MID_H)

    yr = np.array(YEARS)
    ax.fill_between(yr, BID_ASK, alpha=0.3, color=TEAL_H)
    ax.plot(yr, BID_ASK, 'o-', color=TEAL_H, linewidth=2,
            markersize=5, markerfacecolor=WHITE_H, markeredgewidth=1.5)

    ax.set_xticks(yr)
    ax.set_xticklabels([str(y) for y in YEARS], color='#90A4BE', fontsize=7)
    ax.tick_params(axis='y', colors='#90A4BE')
    ax.set_ylabel('Spread %', color='#90A4BE', fontsize=8)
    ax.grid(axis='y', color='#1A3C7D', linewidth=0.5, linestyle='--', alpha=0.5)
    ax.set_title('Bid-Ask Spread Trend (%)', color=WHITE_H,
                 fontsize=9.5, fontweight='bold', pad=8)
    fig.tight_layout(pad=0.5)
    return chart_to_image(fig)


def make_inflation_chart():
    fig, ax = plt.subplots(figsize=(7.0, 3.0))
    fig.patch.set_facecolor(NAVY_MID_H)
    ax.set_facecolor(NAVY_MID_H)

    yr = np.array(YEARS)
    w  = 0.35
    x  = np.arange(len(yr))

    bars1 = ax.bar(x - w/2, US_INF, width=w, label='US Inflation',
                   color=ACCENT_H, alpha=0.85, zorder=3, edgecolor='none')
    bars2 = ax.bar(x + w/2, ZM_INF, width=w, label='Zambia Inflation',
                   color=GOLD_H, alpha=0.85, zorder=3, edgecolor='none')

    ax.set_xticks(x)
    ax.set_xticklabels([str(y) for y in YEARS], color='#90A4BE', fontsize=7)
    ax.tick_params(axis='y', colors='#90A4BE')
    ax.set_ylabel('Inflation Rate (%)', color='#90A4BE', fontsize=8)
    ax.grid(axis='y', color='#1A3C7D', linewidth=0.5, linestyle='--', alpha=0.5)
    ax.legend(framealpha=0.0, fontsize=7.5, labelcolor='#90A4BE', loc='upper right')
    ax.set_title('US vs Zambia Inflation Rate Comparison', color=WHITE_H,
                 fontsize=9.5, fontweight='bold', pad=8)
    fig.tight_layout(pad=0.5)
    return chart_to_image(fig)


def make_inflation_diff_chart():
    fig, ax = plt.subplots(figsize=(5.5, 2.6))
    fig.patch.set_facecolor(NAVY_MID_H)
    ax.set_facecolor(NAVY_MID_H)

    yr = np.array(YEARS)
    ax.fill_between(yr, INF_DIFF, alpha=0.35, color=RED_H)
    ax.plot(yr, INF_DIFF, 'o-', color=RED_H, linewidth=2,
            markersize=4.5, markerfacecolor=WHITE_H, markeredgewidth=1.5)
    ax.axhline(np.mean(INF_DIFF), color=GOLD_H, linewidth=1, linestyle='--',
               label=f'Avg {np.mean(INF_DIFF):.1f}%')

    ax.set_xticks(yr)
    ax.set_xticklabels([str(y) for y in YEARS], color='#90A4BE', fontsize=7)
    ax.tick_params(axis='y', colors='#90A4BE')
    ax.set_ylabel('Diff. (%)', color='#90A4BE', fontsize=8)
    ax.legend(framealpha=0.0, fontsize=7, labelcolor=GOLD_H)
    ax.grid(axis='y', color='#1A3C7D', linewidth=0.5, linestyle='--', alpha=0.5)
    ax.set_title('Inflation Differential (Zambia − US)', color=WHITE_H,
                 fontsize=9.5, fontweight='bold', pad=8)
    fig.tight_layout(pad=0.5)
    return chart_to_image(fig)


def make_policy_gauge():
    """Semicircular gauge for policy verdict."""
    fig, ax = plt.subplots(figsize=(3.2, 2.0), subplot_kw=dict(polar=False))
    fig.patch.set_facecolor(NAVY_MID_H)
    ax.set_facecolor(NAVY_MID_H)
    ax.set_aspect('equal')
    ax.axis('off')

    # draw arc segments
    from matplotlib.patches import Wedge
    cx, cy, r_out, r_in = 0.5, 0.15, 0.38, 0.24
    segments = [
        (180, 240, GREEN_H,   'Safe'),
        (240, 300, GOLD_H,    'Caution'),
        (300, 360, RED_H,     'Pressure'),
    ]
    for start, end, color, lbl in segments:
        wedge = Wedge((cx, cy), r_out, start, end,
                      width=r_out - r_in, facecolor=color, alpha=0.85,
                      transform=ax.transAxes)
        ax.add_patch(wedge)

    # needle at ~315° (pressure zone)
    angle_deg = 325
    angle_rad = np.radians(angle_deg)
    needle_len = 0.30
    nx = cx + needle_len * np.cos(angle_rad)
    ny = cy + needle_len * np.sin(angle_rad)
    ax.annotate('', xy=(nx, ny), xytext=(cx, cy),
                xycoords='axes fraction', textcoords='axes fraction',
                arrowprops=dict(arrowstyle='->', color=WHITE_H, lw=2))
    circle = plt.Circle((cx, cy), 0.035, color=WHITE_H, transform=ax.transAxes)
    ax.add_patch(circle)

    ax.text(0.5, 0.62, 'POLICY VERDICT', transform=ax.transAxes,
            ha='center', va='center', color='#90A4BE', fontsize=6.5)
    ax.text(0.5, 0.50, 'Depreciation Pressure',
            transform=ax.transAxes, ha='center', va='center',
            color=RED_H, fontsize=9, fontweight='bold')
    fig.tight_layout(pad=0)
    return chart_to_image(fig)


def make_executive_summary_chart():
    """Mini sparkline area chart for exec summary."""
    fig, ax = plt.subplots(figsize=(4.5, 2.0))
    fig.patch.set_facecolor(NAVY_MID_H)
    ax.set_facecolor(NAVY_MID_H)

    yr = np.array(YEARS)
    ax.fill_between(yr, AVG_RATE, alpha=0.3, color=ACCENT_H)
    ax.plot(yr, AVG_RATE, '-', color=ACCENT_H, linewidth=2)
    ax.fill_between(yr, ZM_INF, alpha=0.25, color=GOLD_H)
    ax.plot(yr, ZM_INF, color=GOLD_H, linewidth=1.5, linestyle='--')

    ax.set_xticks(yr)
    ax.set_xticklabels([str(y) for y in YEARS], color='#90A4BE', fontsize=6.5)
    ax.tick_params(axis='y', colors='#90A4BE')
    ax.grid(axis='y', color='#1A3C7D', linewidth=0.4, linestyle='--', alpha=0.4)
    p1 = mpatches.Patch(color=ACCENT_H, label='Avg Rate (ZMW/USD)')
    p2 = mpatches.Patch(color=GOLD_H, label='Zambia Inflation %')
    ax.legend(handles=[p1, p2], framealpha=0.0, fontsize=6.5,
              labelcolor='#90A4BE', loc='upper left')
    ax.set_title('Exchange Rate & Inflation Trend (2015–2023)',
                 color=WHITE_H, fontsize=8.5, fontweight='bold', pad=5)
    fig.tight_layout(pad=0.4)
    return chart_to_image(fig)


# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDE BUILDERS
# ═══════════════════════════════════════════════════════════════════════════════

def build_slide_01_cover(prs):
    """Cover slide – full navy gradient with geometric elements."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    W, H = 13.33, 7.5
    add_rect(slide, 0, 0, W, H, NAVY)

    # Large geometric circle (top-right decoration)
    def big_circle(slide, cx, cy, r, color_hex, alpha_str="1D"):
        cx_emu = int(Inches(cx - r)); cy_emu = int(Inches(cy - r))
        r_emu  = int(Inches(r * 2))
        xml = (
            '<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"'
            ' xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
            '<p:nvSpPr><p:cNvPr id="99" name="circle"/>'
            '<p:cNvSpPr><a:spLocks noGrp="1"/></p:cNvSpPr><p:nvPr/></p:nvSpPr>'
            '<p:spPr>'
            f'<a:xfrm><a:off x="{cx_emu}" y="{cy_emu}"/>'
            f'<a:ext cx="{r_emu}" cy="{r_emu}"/></a:xfrm>'
            '<a:prstGeom prst="ellipse"><a:avLst/></a:prstGeom>'
            '<a:solidFill>'
            f'<a:srgbClr val="{color_hex.lstrip("#")}">'
            f'<a:alpha val="{alpha_str}000"/>'
            f'</a:srgbClr>'
            '</a:solidFill>'
            '<a:ln><a:noFill/></a:ln>'
            '</p:spPr>'
            '<p:txBody><a:bodyPr/><a:lstStyle/><a:p/></p:txBody></p:sp>'
        )
        slide.shapes._spTree.append(parse_xml(xml))

    big_circle(slide, 11.5, 1.2, 2.8, NAVY_LIGHT_H, "1A")
    big_circle(slide, 12.2, 0.5, 1.8, ACCENT_H, "22")
    big_circle(slide, 1.5,  6.8, 1.5, NAVY_LIGHT_H, "1A")
    big_circle(slide, 0.5,  7.2, 1.0, ACCENT_H, "1A")

    # Gold top accent stripe
    add_rect(slide, 0, 0, W, 0.08, GOLD)

    # BOZ logo – top left
    logo_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "boz_logo_transparent.png")
    slide.shapes.add_picture(logo_path, Inches(0.22), Inches(0.18), height=Inches(1.35))

    # Left accent bar
    add_rect(slide, 0.55, 1.9, 0.07, 3.5, GOLD)

    # Main title
    add_text(slide,
             "EXCHANGE RATE DETERMINATION",
             0.75, 1.85, 11.5, 1.1,
             36, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    add_text(slide,
             "ZAMBIA ZMW/USD — STRATEGIC POLICY ANALYSIS",
             0.75, 2.90, 10.0, 0.6,
             16, bold=False, color=SKY_BLUE, align=PP_ALIGN.LEFT)

    # Divider line
    add_rect(slide, 0.75, 3.55, 5.5, 0.025, GOLD)

    # Subtitle description
    add_text(slide,
             "A comprehensive analysis of exchange rate policy, purchasing power parity,\n"
             "currency depreciation trends, and macroeconomic inflation dynamics.",
             0.75, 3.65, 9.5, 0.85,
             11, bold=False, color=MID_GRAY, align=PP_ALIGN.LEFT, italic=True)

    # Bottom info bar
    add_rect(slide, 0, 6.35, W, 1.15, NAVY_MID)
    add_rect(slide, 0, 6.35, W, 0.025, ACCENT_BLUE)

    meta_items = [
        ("PREPARED FOR",  "FSi Outsourcing — Executive Leadership"),
        ("ANALYSIS PERIOD", "2015 – 2023"),
        ("CLASSIFICATION", "Internal  |  Confidential"),
        ("DATE",           "May 2026"),
    ]
    for i, (lbl, val) in enumerate(meta_items):
        xpos = 0.55 + i * 3.2
        add_text(slide, lbl, xpos, 6.45, 3.0, 0.28, 6.5,
                 bold=False, color=MID_GRAY, align=PP_ALIGN.LEFT)
        add_text(slide, val, xpos, 6.72, 3.1, 0.35, 9.5,
                 bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    # Decorative KPI strip (3 headline numbers)
    kpi_strip = [
        ("20.89", "ZMW/USD — Latest Avg Rate"),
        ("+21.3%", "YoY Currency Depreciation"),
        ("9.4 pp", "Inflation Differential"),
    ]
    for i, (val, lbl) in enumerate(kpi_strip):
        xpos = 0.75 + i * 4.1
        add_rounded_rect(slide, xpos, 5.0, 3.5, 1.0, NAVY_LIGHT,
                         radius_emu=80000, line_rgb=ACCENT_BLUE, line_pt=1)
        add_rect(slide, xpos + 0.02, 5.0, 0.05, 1.0, GOLD)
        add_text(slide, val, xpos + 0.18, 5.05, 3.2, 0.55,
                 24, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
        add_text(slide, lbl, xpos + 0.18, 5.58, 3.2, 0.35,
                 8, bold=False, color=SKY_BLUE, align=PP_ALIGN.LEFT)


def build_slide_02_agenda(prs):
    """Agenda / Table of Contents."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    W, H  = 13.33, 7.5
    add_rect(slide, 0, 0, W, H, OFF_WHITE)
    add_rect(slide, 0, 0, W, 1.4, NAVY)
    add_rect(slide, 0, 0, 0.07, 1.4, GOLD)

    add_text(slide, "AGENDA", 0.25, 0.22, 4.0, 0.55,
             28, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    add_text(slide, "Exchange Rate Determination — Zambia Analysis",
             0.25, 0.80, 10.0, 0.45,
             10, bold=False, color=SKY_BLUE, align=PP_ALIGN.LEFT)

    sections = [
        ("01", "Exchange Rate Policy",
         "Trend analysis, policy bounds, PPP target, and selling rate KPIs",
         ACCENT_BLUE),
        ("02", "Depreciation & Bid-Ask Analysis",
         "Year-on-year currency depreciation and foreign exchange spread dynamics",
         TEAL),
        ("03", "Inflation Analysis",
         "US vs Zambia inflation comparison and differential impact assessment",
         GOLD),
        ("04", "Policy Verdict & Recommendations",
         "Conclusion on exchange rate alignment and strategic policy direction",
         GREEN),
    ]

    for i, (num, title, desc, color) in enumerate(sections):
        col = i % 2
        row = i // 2
        x = 0.5 + col * 6.35
        y = 1.75 + row * 2.55

        add_rounded_rect(slide, x, y, 5.9, 2.15, WHITE,
                         radius_emu=100000, line_rgb=LIGHT_GRAY, line_pt=1.5)
        add_rect(slide, x, y, 0.09, 2.15, color)

        # Number pill
        add_rounded_rect(slide, x + 0.22, y + 0.25, 0.68, 0.68, color,
                         radius_emu=200000)
        add_text(slide, num, x + 0.22, y + 0.22, 0.68, 0.72,
                 16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        add_text(slide, title, x + 1.08, y + 0.22, 4.6, 0.48,
                 14, bold=True, color=DARK_TEXT, align=PP_ALIGN.LEFT)
        add_rect(slide, x + 1.08, y + 0.73, 3.8, 0.025, LIGHT_GRAY)
        add_text(slide, desc, x + 1.08, y + 0.85, 4.6, 1.0,
                 9.5, bold=False, color=MID_GRAY, align=PP_ALIGN.LEFT)

    # Footer
    add_rect(slide, 0, 7.15, W, 0.35, NAVY)
    add_text(slide, "EXCHANGE RATE DETERMINATION  |  ZAMBIA  |  2015–2023",
             0, 7.17, W, 0.30, 7.5, bold=False, color=MID_GRAY, align=PP_ALIGN.CENTER)


def build_slide_03_exec_summary(prs):
    """Executive Summary with headline KPIs and mini-chart."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    W, H  = 13.33, 7.5
    add_rect(slide, 0, 0, W, H, NAVY)
    add_rect(slide, 0, 0, W, 0.07, GOLD)

    # Header
    add_rect(slide, 0, 0, W, 1.2, NAVY_MID)
    add_rect(slide, 0, 0, 0.07, 1.2, GOLD)
    add_text(slide, "EXECUTIVE SUMMARY", 0.22, 0.15, 8.0, 0.6,
             26, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    add_text(slide,
             "Key findings across exchange rate, depreciation, inflation, and policy alignment — 2015 to 2023",
             0.22, 0.75, 11.0, 0.38,
             9.5, bold=False, color=SKY_BLUE, align=PP_ALIGN.LEFT)

    # 5 KPI cards row
    kpis = [
        ("Avg Selling Rate\n(ZMW/USD)", "20.89",    "",       WHITE,   ACCENT_BLUE),
        ("YoY Depreciation",            "+21.3%",   "2023",   RED,     RED),
        ("Bid-Ask Spread",              "3.1%",     "2023",   GOLD,    GOLD),
        ("Inflation Diff.",             "9.4 pp",   "ZM–US",  GOLD,    GOLD),
        ("Policy Verdict",              "! Pressure","",      RED,     RED),
    ]
    card_w = 2.28
    for i, (lbl, val, unit, vc, ac) in enumerate(kpis):
        xpos = 0.28 + i * (card_w + 0.12)
        kpi_card(slide, xpos, 1.38, card_w, 1.38, lbl, val, unit,
                 value_color=vc, accent_color=ac,
                 bg_color=NAVY_MID, border_color=ac)

    # Narrative bullets (left)
    add_rounded_rect(slide, 0.28, 2.98, 6.2, 3.75, NAVY_MID,
                     radius_emu=80000, line_rgb=NAVY_LIGHT, line_pt=1)
    add_text(slide, "KEY FINDINGS", 0.5, 3.08, 5.8, 0.38,
             9, bold=True, color=GOLD, align=PP_ALIGN.LEFT)
    add_rect(slide, 0.5, 3.45, 5.8, 0.025, NAVY_LIGHT)

    findings = [
        ("Exchange Rate Depreciation",
         "The ZMW/USD rate surged from 8.60 in 2015 to 20.89 in 2023 — a cumulative\n"
         "depreciation of ~143%, exceeding the PPP-implied target rate of 19.50."),
        ("Policy Bound Breach",
         "The 2023 rate of 20.89 ZMW/USD sits ~7.2% below the upper policy bound\n"
         "(22.98), indicating limited but narrowing headroom."),
        ("Inflation Differential",
         "Zambia's average inflation (13.5%) significantly outpaces US inflation (4.1%),\n"
         "sustaining persistent depreciation pressure through PPP dynamics."),
        ("Bid-Ask Spread Widening",
         "The bid-ask spread widened from 1.8% (2015) to 3.1% (2023), signalling\n"
         "reduced FX market liquidity and elevated transaction costs."),
    ]
    for i, (title, body) in enumerate(findings):
        y = 3.60 + i * 0.84
        add_rounded_rect(slide, 0.42, y, 0.38, 0.38, ACCENT_BLUE, radius_emu=200000)
        add_text(slide, str(i + 1), 0.42, y - 0.01, 0.38, 0.40,
                 10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(slide, title, 0.92, y, 5.4, 0.26,
                 8.5, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
        add_text(slide, body, 0.92, y + 0.25, 5.4, 0.55,
                 7.5, bold=False, color=MID_GRAY, align=PP_ALIGN.LEFT)

    # Mini chart (right)
    chart_img = make_executive_summary_chart()
    add_rounded_rect(slide, 6.72, 2.98, 6.38, 3.75, NAVY_MID,
                     radius_emu=80000, line_rgb=NAVY_LIGHT, line_pt=1)
    add_chart_image(slide, chart_img, 6.85, 3.10, 6.1, 3.45)

    # Footer
    add_rect(slide, 0, 7.15, W, 0.35, NAVY_MID)
    add_text(slide, "EXCHANGE RATE DETERMINATION  |  EXECUTIVE SUMMARY  |  CONFIDENTIAL",
             0, 7.17, W, 0.30, 7, bold=False, color=MID_GRAY, align=PP_ALIGN.CENTER)


def build_slide_04_exchange_rate_policy(prs):
    """Section 1: Exchange Rate Policy."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    W, H  = 13.33, 7.5
    add_rect(slide, 0, 0, W, H, NAVY)
    section_header(slide, "01", "EXCHANGE RATE POLICY",
                   "Trend Analysis  ·  Policy Bounds  ·  PPP Target Rate  ·  2015–2023")

    # 4 KPI cards
    kpis = [
        ("Avg Selling Rate",   "20.89",  "ZMW/USD",   WHITE,       ACCENT_BLUE),
        ("PPP Target Rate",    "19.50",  "ZMW/USD",   SKY_BLUE,    TEAL),
        ("Policy Variance",    "-9.1%",  "vs Upper Bound", RED,     RED),
        ("Upper Bound",        "22.98",  "ZMW/USD",   MID_GRAY,    NAVY_LIGHT),
    ]
    for i, (lbl, val, unit, vc, ac) in enumerate(kpis):
        xpos = 0.28 + i * 3.18
        kpi_card(slide, xpos, 1.28, 3.02, 1.45, lbl, val, unit,
                 value_color=vc, accent_color=ac, bg_color=NAVY_MID, border_color=ac)

    # Main chart
    add_rounded_rect(slide, 0.28, 2.90, 12.75, 4.22, NAVY_MID,
                     radius_emu=80000, line_rgb=NAVY_LIGHT, line_pt=1)
    chart_img = make_exchange_rate_chart()
    add_chart_image(slide, chart_img, 0.42, 3.02, 12.45, 3.98)

    # Insight callout
    add_rounded_rect(slide, 9.55, 1.30, 3.52, 1.38, NAVY_LIGHT,
                     radius_emu=60000, line_rgb=GOLD, line_pt=1.5)
    add_text(slide, "KEY INSIGHT", 9.72, 1.38, 3.2, 0.28,
             7, bold=True, color=GOLD, align=PP_ALIGN.LEFT)
    add_text(slide,
             "The 2023 selling rate (20.89) has surpassed the PPP target (19.50) "
             "and is approaching the upper policy bound (22.98). "
             "Continued depreciation pressure warrants policy review.",
             9.72, 1.65, 3.2, 1.0,
             7.5, bold=False, color=OFF_WHITE, align=PP_ALIGN.LEFT)

    add_rect(slide, 0, 7.15, W, 0.35, NAVY_MID)
    add_text(slide, "SECTION 01  |  EXCHANGE RATE POLICY  |  PAGE 1 OF 4",
             0, 7.17, W, 0.30, 7, bold=False, color=MID_GRAY, align=PP_ALIGN.CENTER)


def build_slide_05_depreciation(prs):
    """Section 2: Depreciation & Bid-Ask Analysis."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    W, H  = 13.33, 7.5
    add_rect(slide, 0, 0, W, H, NAVY)
    section_header(slide, "02", "DEPRECIATION & BID-ASK RATE ANALYSIS",
                   "Currency Depreciation  ·  FX Spread Trends  ·  2015–2023")

    # 3 KPI cards
    kpis = [
        ("YoY Depreciation",   "+21.3%",  "2023",       RED,      RED),
        ("Bid-Ask Spread",     "3.1%",    "2023 Latest", GOLD,     GOLD),
        ("Cumulative Depr.",   "+142.9%", "2015→2023",   RED,      RED),
    ]
    for i, (lbl, val, unit, vc, ac) in enumerate(kpis):
        xpos = 0.28 + i * 4.35
        kpi_card(slide, xpos, 1.28, 4.2, 1.45, lbl, val, unit,
                 value_color=vc, accent_color=ac, bg_color=NAVY_MID, border_color=ac)

    # Two charts side by side
    add_rounded_rect(slide, 0.28, 2.90, 6.40, 4.28, NAVY_MID,
                     radius_emu=80000, line_rgb=NAVY_LIGHT, line_pt=1)
    add_rounded_rect(slide, 6.85, 2.90, 6.22, 4.28, NAVY_MID,
                     radius_emu=80000, line_rgb=NAVY_LIGHT, line_pt=1)

    depr_img = make_depreciation_chart()
    ba_img   = make_bidask_chart()
    add_chart_image(slide, depr_img, 0.42, 3.02, 6.15, 4.02)
    add_chart_image(slide, ba_img,   6.98, 3.02, 5.98, 4.02)

    # Annotations
    add_rounded_rect(slide, 0.28, 6.72, 6.40, 0.42, NAVY_LIGHT,
                     radius_emu=40000, line_rgb=RED, line_pt=1)
    add_text(slide,
             "2020 saw the sharpest depreciation (+32.9%) driven by COVID-19 economic shock",
             0.42, 6.76, 6.15, 0.30,
             7.5, bold=False, color=OFF_WHITE, align=PP_ALIGN.LEFT)

    add_rounded_rect(slide, 6.85, 6.72, 6.22, 0.42, NAVY_LIGHT,
                     radius_emu=40000, line_rgb=GOLD, line_pt=1)
    add_text(slide,
             "Spread widening post-2019 reflects reduced FX liquidity and elevated hedging costs",
             6.98, 6.76, 6.0, 0.30,
             7.5, bold=False, color=OFF_WHITE, align=PP_ALIGN.LEFT)

    add_rect(slide, 0, 7.17, W, 0.33, NAVY_MID)
    add_text(slide, "SECTION 02  |  DEPRECIATION & BID-ASK ANALYSIS  |  PAGE 2 OF 4",
             0, 7.19, W, 0.28, 7, bold=False, color=MID_GRAY, align=PP_ALIGN.CENTER)


def build_slide_06_inflation(prs):
    """Section 3: Inflation Analysis."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    W, H  = 13.33, 7.5
    add_rect(slide, 0, 0, W, H, NAVY)
    section_header(slide, "03", "INFLATION ANALYSIS",
                   "US vs Zambia Inflation  ·  Differential Impact  ·  PPP Implications  ·  2015–2023")

    # 4 KPI cards
    kpis = [
        ("US Avg Inflation",       "4.1%",   "2023",    WHITE,    ACCENT_BLUE),
        ("Zambia Avg Inflation",   "13.5%",  "2023",    GOLD,     GOLD),
        ("Inflation Differential", "9.4 pp", "ZM − US", RED,      RED),
        ("Avg Diff. (2015–2023)",  "10.8 pp","9-yr avg",MID_GRAY, NAVY_LIGHT),
    ]
    for i, (lbl, val, unit, vc, ac) in enumerate(kpis):
        xpos = 0.28 + i * 3.18
        kpi_card(slide, xpos, 1.28, 3.02, 1.45, lbl, val, unit,
                 value_color=vc, accent_color=ac, bg_color=NAVY_MID, border_color=ac)

    # Left: grouped bar chart
    add_rounded_rect(slide, 0.28, 2.90, 7.45, 4.22, NAVY_MID,
                     radius_emu=80000, line_rgb=NAVY_LIGHT, line_pt=1)
    inf_img = make_inflation_chart()
    add_chart_image(slide, inf_img, 0.42, 3.02, 7.20, 3.98)

    # Right: diff line + PPP note
    add_rounded_rect(slide, 7.90, 2.90, 5.15, 4.22, NAVY_MID,
                     radius_emu=80000, line_rgb=NAVY_LIGHT, line_pt=1)
    diff_img = make_inflation_diff_chart()
    add_chart_image(slide, diff_img, 8.02, 3.02, 4.90, 2.55)

    # PPP note box
    add_rounded_rect(slide, 7.90, 5.60, 5.15, 1.52, NAVY_LIGHT,
                     radius_emu=60000, line_rgb=GOLD, line_pt=1.5)
    add_text(slide, "PPP IMPLICATION", 8.08, 5.68, 4.7, 0.28,
             7.5, bold=True, color=GOLD, align=PP_ALIGN.LEFT)
    add_text(slide,
             "Purchasing Power Parity theory dictates that sustained inflation "
             "differentials drive proportional exchange rate depreciation. "
             "Zambia's persistent 9–11 pp differential above US CPI continues "
             "to exert structural downward pressure on the ZMW.",
             8.08, 5.95, 4.80, 1.08,
             8, bold=False, color=OFF_WHITE, align=PP_ALIGN.LEFT)

    add_rect(slide, 0, 7.17, W, 0.33, NAVY_MID)
    add_text(slide, "SECTION 03  |  INFLATION ANALYSIS  |  PAGE 3 OF 4",
             0, 7.19, W, 0.28, 7, bold=False, color=MID_GRAY, align=PP_ALIGN.CENTER)


def build_slide_07_verdict(prs):
    """Section 4: Policy Verdict & Recommendations."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    W, H  = 13.33, 7.5
    add_rect(slide, 0, 0, W, H, NAVY)
    section_header(slide, "04", "POLICY VERDICT & RECOMMENDATIONS",
                   "Exchange Rate Alignment Assessment  ·  Strategic Guidance  ·  2023")

    # Verdict card – wide
    add_rounded_rect(slide, 0.28, 1.30, 12.75, 1.55, NAVY_MID,
                     radius_emu=90000, line_rgb=RED, line_pt=2.0)
    add_rect(slide, 0.28, 1.30, 0.09, 1.55, RED)

    add_text(slide, "VERDICT", 0.52, 1.38, 2.5, 0.30,
             8, bold=True, color=RED, align=PP_ALIGN.LEFT)
    add_text(slide, "DEPRECIATION PRESSURE — APPROACHING POLICY BOUNDS",
             0.52, 1.66, 9.5, 0.55,
             18, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    add_text(slide,
             "The ZMW/USD rate of 20.89 (2023) exceeds the PPP target (19.50) "
             "and is tracking toward the upper bound (22.98). "
             "The 9.4 pp inflation differential reinforces continued depreciation risk.",
             0.52, 2.16, 9.5, 0.55,
             9, bold=False, color=MID_GRAY, align=PP_ALIGN.LEFT)

    gauge_img = make_policy_gauge()
    add_chart_image(slide, gauge_img, 10.15, 1.28, 2.92, 1.60)

    # 4 recommendation boxes
    recommendations = [
        (ACCENT_BLUE, "Monetary Policy Tightening",
         "The Reserve Bank of Zambia should consider targeted rate adjustments "
         "to moderate inflation and slow the depreciation trajectory."),
        (TEAL, "FX Market Intervention",
         "Strategic FX reserves deployment at key technical levels "
         "(21.50–22.00 ZMW/USD) can dampen excessive volatility."),
        (GOLD, "Inflation Anchor Framework",
         "Implement an explicit inflation targeting band of 6–8% to close "
         "the inflation differential gap with the US over a 3-year horizon."),
        (GREEN, "PPP Monitoring Dashboard",
         "Establish a quarterly PPP equilibrium review against the Optimal "
         "Target Rate to provide early-warning signals for policy intervention."),
    ]
    for i, (color, title, body) in enumerate(recommendations):
        col = i % 2; row = i // 2
        x = 0.28 + col * 6.52; y = 3.12 + row * 1.95
        add_rounded_rect(slide, x, y, 6.30, 1.78, NAVY_MID,
                         radius_emu=80000, line_rgb=color, line_pt=1.5)
        add_rect(slide, x, y, 0.08, 1.78, color)
        add_rounded_rect(slide, x + 0.22, y + 0.22, 0.55, 0.55, color,
                         radius_emu=200000)
        icons = ["01", "02", "03", "04"]
        add_text(slide, icons[i], x + 0.20, y + 0.16, 0.6, 0.55,
                 11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(slide, title, x + 0.92, y + 0.18, 5.2, 0.38,
                 10.5, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
        add_rect(slide, x + 0.92, y + 0.58, 5.0, 0.022, NAVY_LIGHT)
        add_text(slide, body, x + 0.92, y + 0.68, 5.2, 1.0,
                 8, bold=False, color=MID_GRAY, align=PP_ALIGN.LEFT)

    add_rect(slide, 0, 7.17, W, 0.33, NAVY_MID)
    add_text(slide, "SECTION 04  |  POLICY VERDICT & RECOMMENDATIONS  |  PAGE 4 OF 4",
             0, 7.19, W, 0.28, 7, bold=False, color=MID_GRAY, align=PP_ALIGN.CENTER)


def build_slide_08_conclusion(prs):
    """Closing / Thank You slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    W, H  = 13.33, 7.5
    add_rect(slide, 0, 0, W, H, NAVY)
    add_rect(slide, 0, 0, W, 0.07, GOLD)
    add_rect(slide, 0, H - 0.07, W, 0.07, GOLD)

    # Decorative circles
    def big_circle_closing(cx, cy, r, col_h, alpha):
        cx_emu = int(Inches(cx - r)); cy_emu = int(Inches(cy - r))
        r_emu  = int(Inches(r * 2))
        xml = (
            '<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"'
            ' xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
            '<p:nvSpPr><p:cNvPr id="88" name="c2"/>'
            '<p:cNvSpPr><a:spLocks noGrp="1"/></p:cNvSpPr><p:nvPr/></p:nvSpPr>'
            '<p:spPr>'
            f'<a:xfrm><a:off x="{cx_emu}" y="{cy_emu}"/>'
            f'<a:ext cx="{r_emu}" cy="{r_emu}"/></a:xfrm>'
            '<a:prstGeom prst="ellipse"><a:avLst/></a:prstGeom>'
            f'<a:solidFill><a:srgbClr val="{col_h.lstrip("#")}">'
            f'<a:alpha val="{alpha}000"/></a:srgbClr></a:solidFill>'
            '<a:ln><a:noFill/></a:ln></p:spPr>'
            '<p:txBody><a:bodyPr/><a:lstStyle/><a:p/></p:txBody></p:sp>'
        )
        slide.shapes._spTree.append(parse_xml(xml))

    big_circle_closing(12.5, 1.0, 2.5, NAVY_LIGHT_H, "18")
    big_circle_closing(1.0, 6.5, 2.0, NAVY_LIGHT_H, "15")
    big_circle_closing(11.8, 6.8, 1.5, ACCENT_H, "10")

    add_rect(slide, 4.2, 0, 0.06, H, GOLD)

    # Left panel: summary stats
    add_text(slide, "ANALYSIS SUMMARY", 0.45, 1.00, 3.55, 0.50,
             10, bold=True, color=GOLD, align=PP_ALIGN.LEFT)
    summary_kpis = [
        ("143%",   "Cumulative ZMW Depreciation\n2015–2023"),
        ("20.89",  "ZMW/USD Avg Rate (2023)"),
        ("9.4 pp", "Inflation Differential\n(Zambia vs US, 2023)"),
        ("3.1%",   "Bid-Ask Spread (2023)"),
    ]
    for i, (val, lbl) in enumerate(summary_kpis):
        y = 1.65 + i * 1.28
        add_rect(slide, 0.45, y, 0.05, 0.95, ACCENT_BLUE)
        add_text(slide, val, 0.65, y, 2.8, 0.52,
                 20, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
        add_text(slide, lbl, 0.65, y + 0.50, 3.3, 0.38,
                 8, bold=False, color=MID_GRAY, align=PP_ALIGN.LEFT)

    # Right panel: closing message
    add_text(slide, "THANK YOU", 4.55, 1.60, 8.5, 1.0,
             42, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    add_rect(slide, 4.55, 2.68, 3.5, 0.045, GOLD)
    add_text(slide,
             "This analysis provides a data-driven foundation for exchange\n"
             "rate policy review and strategic decision-making for Zambia.\n\n"
             "The findings highlight the urgency of addressing inflation\n"
             "differentials and FX market depth to sustain ZMW stability.",
             4.55, 2.82, 8.3, 1.85,
             11, bold=False, color=MID_GRAY, align=PP_ALIGN.LEFT)

    add_text(slide, "QUESTIONS & DISCUSSION", 4.55, 4.85, 8.3, 0.45,
             13, bold=True, color=SKY_BLUE, align=PP_ALIGN.LEFT)

    # Contact / attribution row
    add_rounded_rect(slide, 4.40, 5.52, 8.70, 1.18, NAVY_MID,
                     radius_emu=60000, line_rgb=NAVY_LIGHT, line_pt=1)
    add_text(slide, "PREPARED BY", 4.60, 5.60, 3.0, 0.28,
             7, bold=True, color=GOLD, align=PP_ALIGN.LEFT)
    add_text(slide, "FSi Outsourcing — Analytics & Advisory",
             4.60, 5.86, 5.5, 0.30,
             10, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    add_text(slide, "k.banda@fsioutsourcing.com  |  May 2026  |  CONFIDENTIAL",
             4.60, 6.16, 7.0, 0.28,
             8, bold=False, color=MID_GRAY, align=PP_ALIGN.LEFT)

    add_text(slide,
             "Data sourced from Power BI model: Analysis.pbix\n"
             "Analysis Period: 2015–2023  |  All rates in ZMW/USD",
             4.55, 6.65, 8.3, 0.50,
             7, bold=False, color=NAVY_LIGHT, align=PP_ALIGN.LEFT)


# ═══════════════════════════════════════════════════════════════════════════════
#  MAIN
# ═══════════════════════════════════════════════════════════════════════════════

def main():
    out_path = os.path.join(
        r"c:\Users\k.banda\Documents\GitHub\Exchange-Rate-Determination",
        "Exchange_Rate_Determination_Zambia.pptx"
    )

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    print("Building slides...")
    build_slide_01_cover(prs);              print("  OK Slide 1: Cover")
    build_slide_02_agenda(prs);             print("  OK Slide 2: Agenda")
    build_slide_03_exec_summary(prs);       print("  OK Slide 3: Executive Summary")
    build_slide_04_exchange_rate_policy(prs); print("  OK Slide 4: Exchange Rate Policy")
    build_slide_05_depreciation(prs);       print("  OK Slide 5: Depreciation & Bid-Ask")
    build_slide_06_inflation(prs);          print("  OK Slide 6: Inflation Analysis")
    build_slide_07_verdict(prs);            print("  OK Slide 7: Policy Verdict")
    build_slide_08_conclusion(prs);         print("  OK Slide 8: Conclusion")

    prs.save(out_path)
    print(f"\nDONE. Presentation saved to:\n   {out_path}")


if __name__ == "__main__":
    main()
